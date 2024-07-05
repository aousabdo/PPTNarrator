import os
import json
import logging
from pdf2image import convert_from_path
from pptx import Presentation
import subprocess
from pptx.util import Inches
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from claude_narrator import get_narration_from_claude
from claude_narrator import get_summary_from_claude

def ppt_to_png(ppt_path, dpi=300):
    # Get the base name of the input file (without extension)
    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    
    # Create a new output folder based on the PPT file name with '_images' tag
    output_folder = os.path.join(os.path.dirname(ppt_path), f"{base_name}_images")
    os.makedirs(output_folder, exist_ok=True)

    # Construct the LibreOffice command to convert PPT to PDF
    libreoffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    pdf_path = os.path.join(output_folder, f"{base_name}.pdf")
    command = [
        libreoffice_path,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_folder,
        ppt_path
    ]

    # Run the command to convert PPT to PDF
    try:
        subprocess.run(command, check=True, capture_output=True, text=True)
        print(f"Converted {ppt_path} to PDF")
    except subprocess.CalledProcessError as e:
        print(f"An error occurred during conversion to PDF: {e}")
        print(f"Error output: {e.output}")
        return

    # Convert PDF to high-resolution PNG
    try:
        images = convert_from_path(pdf_path, dpi=dpi)
        for i, image in enumerate(images):
            image.save(os.path.join(output_folder, f"slide_{i+1:03d}.png"), "PNG")
        print(f"Converted PDF to {len(images)} high-resolution PNG images")
    except Exception as e:
        print(f"An error occurred during conversion from PDF to PNG: {e}")

    # Clean up the temporary PDF file
    os.remove(pdf_path)

    print(f"All high-resolution images have been saved to: {output_folder}")
    return output_folder

def generate_narration(image_path, slide_number, total_slides, presentation_summary, output_dir):
    narration = get_narration_from_claude(image_path, slide_number, total_slides, presentation_summary)
    if narration:
        text_output_path = os.path.join(output_dir, f"slide_{slide_number:03d}_narration.txt")
        with open(text_output_path, 'w', encoding='utf-8') as f:
            f.write(narration)
        return narration
    else:
        return f"Unable to generate narration for slide {slide_number}."

def generate_output_path(slide_number, output_dir):
    return os.path.join(output_dir, f"slide_{slide_number:03d}.mp3")

def process_slides(ppt_path, output_dir):
    # Get the presentation summary first
    presentation_summary = get_presentation_summary(ppt_path, output_dir)
    
    # Get the base name of the input file (without extension)
    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    
    # Create a new output folder based on the PPT file name with '_images' tag
    images_folder = os.path.join(output_dir, f"{base_name}_images")
    
    # Convert PPT to PNG images
    images_folder = ppt_to_png(ppt_path, dpi=300)
    
    # Get all PNG files in the images folder
    image_paths = sorted([os.path.join(images_folder, f) for f in os.listdir(images_folder) if f.endswith('.png')])
    
    narrations = []
    for i, image_path in enumerate(image_paths):
        slide_number = i + 1
        total_slides = len(image_paths)
        text_to_speak = generate_narration(image_path, slide_number, total_slides, presentation_summary, output_dir)
        output_path = generate_output_path(slide_number, output_dir)
        narrations.append((text_to_speak, output_path))
    return narrations

def add_audio_to_ppt(ppt_path, audio_files):
    prs = Presentation(ppt_path)
    
    for i, slide in enumerate(prs.slides):
        if i < len(audio_files):
            audio_path = audio_files[i]
            left = top = Inches(0)
            height = width = Inches(1)  # Small, but not invisible
            
            # Add the audio file to the slide
            audio = slide.shapes.add_movie(audio_path, left, top, width, height)
            
            # Get the XML element
            element = audio.element
            
            # Define namespaces
            namespaces = {
                'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
                'p': "http://schemas.openxmlformats.org/presentationml/2006/main",
                'r': "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                'p14': "http://schemas.microsoft.com/office/powerpoint/2010/main"
            }

            # Find or create nvPr element
            nvPr = element.find('.//p:nvPr', namespaces)
            if nvPr is None:
                nvPicPr = element.find('.//p:nvPicPr', namespaces)
                nvPr = etree.SubElement(nvPicPr, '{%s}nvPr' % namespaces['p'])

            # Add audioFile element
            audioFile = nvPr.find('a:audioFile', namespaces)
            if audioFile is None:
                audioFile = etree.SubElement(nvPr, '{%s}audioFile' % namespaces['a'])
            audioFile.set('embed', 'rId1')
            audioFile.set('name', '')

            # Add extLst element
            extLst = nvPr.find('p:extLst', namespaces)
            if extLst is None:
                extLst = etree.SubElement(nvPr, '{%s}extLst' % namespaces['p'])

            # Add ext element
            ext = extLst.find('p:ext', namespaces)
            if ext is None:
                ext = etree.SubElement(extLst, '{%s}ext' % namespaces['p'])
            ext.set('uri', '{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}')

            # Add media element
            media = ext.find('p14:media', namespaces)
            if media is None:
                media = etree.SubElement(ext, '{%s}media' % namespaces['p14'])

            # Add trim element
            trim = media.find('p14:trim', namespaces)
            if trim is None:
                trim = etree.SubElement(media, '{%s}trim' % namespaces['p14'])
            trim.set('st', '0')

            # Add play element to set autoplay
            play = media.find('p14:play', namespaces)
            if play is None:
                play = etree.SubElement(media, '{%s}play' % namespaces['p14'])
            play.set('auto', '1')

    # Save the modified presentation
    output_path = os.path.join(os.path.dirname(ppt_path), os.path.basename(ppt_path).replace('.pptx', '_with_audio.pptx'))
    prs.save(output_path)
    print(f"Presentation with audio saved as: {output_path}")
    return output_path

def extract_presentation_text(ppt_path):
    prs = Presentation(ppt_path)
    full_text = []
    for slide_number, slide in enumerate(prs.slides, 1):
        slide_text = [f"Slide {slide_number}:"]
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        
        # Extract speaker notes
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = notes_slide.notes_text_frame.text
            if notes_text:
                slide_text.append(f"Speaker Notes: {notes_text}")
        
        full_text.append('\n'.join(slide_text))
    
    return '\n\n'.join(full_text)

def get_presentation_summary(ppt_path, output_dir):
    # Get the base name of the presentation file (without extension)
    ppt_base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    
    # Create a more descriptive summary file name
    summary_file = os.path.join(output_dir, f"{ppt_base_name}_presentation_summary.json")
    
    # Check if summary already exists
    if os.path.exists(summary_file):
        try:
            with open(summary_file, 'r') as f:
                summary = json.load(f)
            if summary:  # Check if the loaded summary is not empty
                return summary
            else:
                logging.warning(f"Existing summary file {summary_file} is empty. Generating new summary.")
        except json.JSONDecodeError:
            logging.warning(f"Error decoding existing summary file {summary_file}. Generating new summary.")
    else:
        logging.info(f"Summary file {summary_file} not found. Generating new summary.")
    
    # Generate new summary
    full_text = extract_presentation_text(ppt_path)
    summary = get_summary_from_claude(full_text)
    
    # Save summary
    with open(summary_file, 'w') as f:
        json.dump(summary, f)
    
    return summary